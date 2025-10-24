"""
PPTX to Markdown Converter
Converts Microsoft PowerPoint presentations to markdown outline format
"""
import os
from pathlib import Path
from datetime import datetime, timezone
from typing import Dict, Any
from pptx import Presentation


def convert(file_path: Path, output_dir: Path) -> Dict[str, Any]:
    """
    Convert PPTX file to markdown outline format
    
    Args:
        file_path: Path to the input PPTX file
        output_dir: Directory where output should be written
    
    Returns:
        Dictionary with conversion result
    """
    try:
        # Create pptx output directory if it doesn't exist
        pptx_dir = output_dir / "pptx"
        pptx_dir.mkdir(parents=True, exist_ok=True)        
        # Load the presentation
        prs = Presentation(str(file_path))
        
        # Build markdown content
        markdown_lines = []
        
        # Add metadata header
        markdown_lines.extend([
            "---",
            f"source: {file_path.relative_to(file_path.parent)}",
            f"converted_at: {datetime.utcnow().isoformat()}Z",
                f"converted_at: {datetime.now(timezone.utc).isoformat()}Z",
            "converter: pptx_to_md",
            "---",
            ""
        ])
        
        # Process each slide
        for slide_idx, slide in enumerate(prs.slides, 1):
            # Get slide title (usually from the first shape)
            slide_title = f"Slide {slide_idx}"
            title_found = False
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Often the title is in the first shape or has a title-like style
                    if shape.has_text_frame:
                        text_frame = shape.text_frame
                        if text_frame.text.strip():
                            # Check if this looks like a title (short text, often first shape)
                            if len(text_frame.text.strip()) < 100 and not title_found:
                                slide_title = text_frame.text.strip().replace('\n', ' ').replace('\r', ' ')
                                title_found = True
                                break
            
            # Add slide header
            markdown_lines.append(f"## {slide_title}")
            
            # Collect all text from shapes
            slide_content = []
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != slide_title:
                            slide_content.append(text)  # Add slide content as bullet points
            for content in slide_content:
                # Split by lines and add as separate bullet points
                for line in content.split('\n'):
                    line = line.strip()
                    if line:
                        markdown_lines.append(f"- {line}")
            
            # Add speaker notes if present
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        markdown_lines.append("\n**Notes:**")
                        for note_line in notes_text.split('\n'):
                            note_line = note_line.strip()
                            if note_line:
                                markdown_lines.append(f"- {note_line}")
            
            # Add blank line between slides
            markdown_lines.append("")
        
        # Join all lines
        markdown_content = "\n".join(markdown_lines).strip()
        
        # Write output file
        output_filename = file_path.stem + ".md"
        output_file = pptx_dir / output_filename
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        
        return {
            'status': 'success',
            'output': str(output_file.relative_to(output_dir.parent)),
            'converter': 'pptx_to_md',
            'metadata': {
                'slides': len(prs.slides),
                'has_notes': any(slide.has_notes_slide for slide in prs.slides)
            }
        }
        
    except Exception as e:
        return {
            'status': 'failed',
            'error': str(e),
            'converter': 'pptx_to_md'
        }
